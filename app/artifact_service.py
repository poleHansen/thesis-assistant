from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path
import re
import shutil
import subprocess
import textwrap
import zipfile

from app.domain import ArtifactBundle, PaperDocument, PaperNode, ProjectState, TemplateManifest
from app.storage import ProjectStorage
from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
from docx.shared import Pt  # type: ignore


DOCX_PLACEHOLDER_PATTERN = re.compile(r"\{\{\s*([\w\-.\u4e00-\u9fff ]+)\s*\}\}")

SECTION_ALIAS_GROUPS = [
    ["引言", "绪论", "第1章 绪论", "第一章 绪论"],
    ["方法", "方法设计", "方法设计与实现", "第3章 方法设计与实现", "第三章 方法设计与实现"],
    ["实验", "实验结果与分析", "实验设计", "第4章 实验结果与分析", "第四章 实验结果与分析"],
    ["结论", "总结与展望", "结论与展望", "第5章 总结与展望", "第五章 总结与展望"],
    ["摘要", "中文摘要"],
    ["参考文献"],
]


@dataclass(frozen=True, slots=True)
class ThesisHtmlStyle:
    font_family: str
    font_size_pt: int
    bold: bool = False
    text_align: str = "justify"
    first_line_indent_chars: int = 0
    margin_top_pt: int = 0
    margin_bottom_pt: int = 12
    line_height: float = 1.75


@dataclass(frozen=True, slots=True)
class ThesisHtmlStyleProfile:
    title: ThesisHtmlStyle
    author_meta: ThesisHtmlStyle
    cover_label: ThesisHtmlStyle
    heading_1: ThesisHtmlStyle
    heading_2: ThesisHtmlStyle
    body: ThesisHtmlStyle


DEFAULT_THESIS_HTML_STYLE_PROFILE = ThesisHtmlStyleProfile(
    title=ThesisHtmlStyle(
        font_family="宋体",
        font_size_pt=16,
        bold=True,
        text_align="center",
        margin_top_pt=0,
        margin_bottom_pt=64,
        line_height=1.5,
    ),
    author_meta=ThesisHtmlStyle(font_family="宋体", font_size_pt=16, text_align="center"),
    cover_label=ThesisHtmlStyle(font_family="宋体", font_size_pt=12, text_align="left"),
    heading_1=ThesisHtmlStyle(
        font_family="宋体",
        font_size_pt=15,
        bold=True,
        margin_top_pt=18,
        margin_bottom_pt=6,
    ),
    heading_2=ThesisHtmlStyle(
        font_family="宋体",
        font_size_pt=14,
        bold=True,
        margin_top_pt=12,
        margin_bottom_pt=6,
    ),
    body=ThesisHtmlStyle(
        font_family="宋体",
        font_size_pt=12,
        first_line_indent_chars=2,
        margin_top_pt=0,
        margin_bottom_pt=6,
        line_height=1.0,
    ),
)


class ArtifactService:
    def __init__(self, storage: ProjectStorage) -> None:
        self.storage = storage
        self.thesis_html_style_profile = DEFAULT_THESIS_HTML_STYLE_PROFILE

    def render_all(self, state: ProjectState) -> ProjectState:
        artifacts_dir = self.storage.ensure_project_tree(state.project_id) / "artifacts"
        reports_dir = artifacts_dir / "reports"
        code_dir = artifacts_dir / "code"
        slides_dir = artifacts_dir / "slides"
        reports_dir.mkdir(parents=True, exist_ok=True)
        code_dir.mkdir(parents=True, exist_ok=True)
        slides_dir.mkdir(parents=True, exist_ok=True)

        word_template_path = state.template_source.template_path if state.template_source else None
        ppt_template_path = self._resolve_ppt_template_path(state)

        state.artifacts = ArtifactBundle(
            literature_review=self._write_literature_review(reports_dir, state),
            innovation_report=self._write_text(
                reports_dir / "innovation_report.md",
                self._innovation_report(state),
            ),
            experiment_plan=self._write_docx_like(
                reports_dir / "experiment_plan.docx",
                self._experiment_plan_text(state),
                manifest=state.template_manifest,
            ),
            procedure=self._write_docx_like(
                reports_dir / "procedure.docx",
                str(state.result_schema.get("procedure_document", "")),
                manifest=state.template_manifest,
            ),
            thesis_html=self._write_text(
                reports_dir / "thesis_source.html",
                self._render_thesis_html(state),
            ),
            thesis_docx=self._write_thesis_docx(
                reports_dir / "thesis.docx",
                state,
                template_path=word_template_path,
                manifest=state.template_manifest,
            ),
            thesis_pdf=self._write_minimal_pdf(
                reports_dir / "thesis.pdf",
                self._thesis_text(state),
            ),
            code_zip=self._write_code_zip(code_dir / "code_bundle.zip", state),
            defense_pptx=self._write_pptx_like(
                slides_dir / "defense.pptx",
                self._ppt_text(state),
                template_path=ppt_template_path,
                manifest=state.template_manifest,
            ),
            qa_report=self._write_text(
                reports_dir / "qa_report.json",
                json.dumps(
                    {"findings": state.review_findings, "warnings": state.warnings},
                    ensure_ascii=False,
                    indent=2,
                ),
            ),
        )
        return state

    def _resolve_ppt_template_path(self, state: ProjectState) -> str | None:
        uploaded_path = state.result_schema.get("ppt_template_path")
        if isinstance(uploaded_path, str) and uploaded_path.strip():
            return uploaded_path
        if state.template_source and state.template_source.ppt_template_path:
            return state.template_source.ppt_template_path
        return None

    def _render_thesis_html(self, state: ProjectState) -> str:
        profile = self.thesis_html_style_profile
        paper_document = state.paper_document or self._paper_document_from_sections(state)
        title = self._escape_html((paper_document.title if paper_document else state.request.topic) or "毕业论文")
        meta_lines = self._build_cover_meta_lines(state)
        html_parts = [
            "<html>",
            "<head><meta http-equiv=\"Content-Type\" content=\"text/html; charset=utf-8\"></head>",
            '<body style="margin:0;padding:0;background:#ffffff;">',
            '<div style="width:720px;min-height:960px;margin:0 auto;padding:72pt 64pt 48pt 64pt;box-sizing:border-box;">',
            '<div style="height:96pt;"></div>',
            f'<p style="{self._style_to_inline_css(profile.title)};margin:0 0 64pt 0;letter-spacing:1pt;">{title}</p>',
        ]
        for line in meta_lines:
            label, _, value = line.partition("：")
            html_parts.append(
                '<p '
                f'title="{self._escape_html(line)}" '
                'style="'
                f'{self._style_to_inline_css(profile.cover_label)}'
                ';display:flex;align-items:flex-end;gap:12pt;margin:0 0 18pt 0;">'
                f'<span style="display:inline-block;min-width:84pt;">{self._escape_html(label)}：</span>'
                '<span style="flex:1;border-bottom:1px solid #222;padding:0 0 2pt 0;min-height:18pt;">'
                f'{self._escape_html(value)}'
                '</span>'
                '</p>'
            )
        html_parts.append('<div style="height:48pt;"></div>')
        if paper_document:
            for node in paper_document.nodes:
                html_parts.extend(self._render_paper_node_html(node, state))
        else:
            for section in state.paper_outline:
                html_parts.append(f'<p style="{self._style_to_inline_css(profile.heading_1)}">{self._escape_html(section)}</p>')
                body_text = self._resolve_section_text(section, state)
                for paragraph in self._split_html_paragraphs(body_text):
                    html_parts.append(f'<p style="{self._style_to_inline_css(profile.body)}">{self._escape_html(paragraph)}</p>')
                if "实验" in section:
                    html_parts.extend(self._render_result_analysis_html(state))
        html_parts.append("</div>")
        html_parts.append("</body>")
        html_parts.append("</html>")
        return "\n".join(html_parts)

    def _render_paper_node_html(self, node: PaperNode, state: ProjectState) -> list[str]:
        profile = self.thesis_html_style_profile
        parts = [f'<p style="{self._style_to_inline_css(profile.heading_1)}">{self._escape_html(node.title)}</p>']
        for paragraph in node.paragraphs:
            for block in self._split_html_paragraphs(paragraph):
                parts.append(f'<p style="{self._style_to_inline_css(profile.body)}">{self._escape_html(block)}</p>')
        for child in node.children:
            parts.append(f'<p style="{self._style_to_inline_css(profile.heading_2)}">{self._escape_html(child.title)}</p>')
            for paragraph in child.paragraphs:
                for block in self._split_html_paragraphs(paragraph):
                    parts.append(f'<p style="{self._style_to_inline_css(profile.body)}">{self._escape_html(block)}</p>')
        if "实验" in node.title:
            parts.extend(self._render_result_analysis_html(state))
        return parts

    def _render_result_analysis_html(self, state: ProjectState) -> list[str]:
        profile = self.thesis_html_style_profile
        parts: list[str] = []
        for block in self._format_result_analysis_blocks(state):
            if block.startswith("### "):
                parts.append(
                    f'<p style="{self._style_to_inline_css(profile.heading_2)}">{self._escape_html(block[4:])}</p>'
                )
                continue
            if block.strip():
                parts.append(f'<p style="{self._style_to_inline_css(profile.body)}">{self._escape_html(block)}</p>')
        return parts

    def _build_cover_meta_lines(self, state: ProjectState) -> list[str]:
        manifest = state.template_manifest or TemplateManifest([], {}, [], [], [], "", {}, {})
        template_name = state.template_source.template_name if state.template_source else "未设置"
        placeholder_values = self._thesis_placeholder_values(state)
        lines = [
            f"课题名称：{state.request.topic}",
            f"论文模板：{template_name}",
        ]
        preferred_fields = [field for field in manifest.cover_fields if field and field not in {"题目", "课题", "标题"}]
        if not preferred_fields:
            preferred_fields = ["学校", "学院", "专业", "学生姓名", "学号", "指导教师", "完成日期"]
        for field in preferred_fields:
            value = str(placeholder_values.get(f"cover.{field}", "")).strip() or "待填写"
            lines.append(f"{field}：{value}")
        return lines

    def _style_to_inline_css(self, style: ThesisHtmlStyle) -> str:
        css_parts = [
            f"font-family:{style.font_family}",
            f"font-size:{style.font_size_pt}pt",
            f"font-weight:{'bold' if style.bold else 'normal'}",
            f"text-align:{style.text_align}",
            f"margin:{style.margin_top_pt}pt 0 {style.margin_bottom_pt}pt 0",
            f"line-height:{style.line_height}",
        ]
        if style.first_line_indent_chars > 0:
            css_parts.append(f"text-indent:{style.first_line_indent_chars}em")
        return ";".join(css_parts)

    def _split_html_paragraphs(self, text: str) -> list[str]:
        return [segment.strip() for segment in re.split(r"\n\s*\n", text) if segment.strip()]

    def _escape_html(self, text: str) -> str:
        return (
            str(text)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;")
        )

    def _convert_html_to_docx(self, html_path: Path, docx_path: Path) -> bool:
        pandoc_path = shutil.which("pandoc")
        if not pandoc_path:
            return False
        command = [pandoc_path, str(html_path), "-o", str(docx_path)]
        completed = subprocess.run(command, capture_output=True, text=True, check=False)
        return completed.returncode == 0 and docx_path.exists()

    def _write_literature_review(self, reports_dir: Path, state: ProjectState) -> str:
        try:
            from openpyxl import Workbook  # type: ignore

            path = reports_dir / "literature_review.xlsx"
            workbook = Workbook()
            sheet = workbook.active
            sheet.title = "literature"
            headers = [
                "title",
                "problem",
                "method",
                "dataset",
                "metrics",
                "conclusion",
                "limitations",
                "source",
                "doi_or_url",
                "evidence_source",
                "confidence",
                "evidence_quote",
                "pdf_path",
                "pdf_parse_status",
                "pdf_parse_message",
                "citation_count",
                "is_fallback",
                "needs_review",
                "review_note",
            ]
            sheet.append(headers)
            for row in state.survey_table:
                sheet.append([row.get(header, "") for header in headers])
            workbook.save(path)
            return str(path)
        except Exception:
            path = reports_dir / "literature_review.csv"
            headers = [
                "title",
                "problem",
                "method",
                "dataset",
                "metrics",
                "conclusion",
                "limitations",
                "source",
                "doi_or_url",
                "evidence_source",
                "confidence",
                "evidence_quote",
                "pdf_path",
                "pdf_parse_status",
                "pdf_parse_message",
                "citation_count",
                "is_fallback",
                "needs_review",
                "review_note",
            ]
            with path.open("w", encoding="utf-8-sig", newline="") as handle:
                writer = csv.DictWriter(handle, fieldnames=headers)
                writer.writeheader()
                writer.writerows(state.survey_table)
            return str(path)

    def _innovation_report(self, state: ProjectState) -> str:
        selected = state.selected_innovation
        summary = state.result_schema.get("gap_analysis", {}) if isinstance(state.result_schema, dict) else {}
        recommendation = (
            state.result_schema.get("innovation_recommendation", {}) if isinstance(state.result_schema, dict) else {}
        )
        mode = (
            summary.get("mode", selected.evidence_mode if selected else "fallback")
            if isinstance(summary, dict)
            else "fallback"
        )
        common_methods = "、".join(summary.get("common_methods", [])[:3]) if isinstance(summary, dict) else ""
        common_datasets = "、".join(summary.get("common_datasets", [])[:3]) if isinstance(summary, dict) else ""
        common_metrics = "、".join(summary.get("common_metrics", [])[:3]) if isinstance(summary, dict) else ""
        common_limitations = "、".join(summary.get("common_limitations", [])[:3]) if isinstance(summary, dict) else ""
        rare_methods = "、".join(summary.get("rare_methods", [])[:3]) if isinstance(summary, dict) else ""
        rare_datasets = "、".join(summary.get("rare_datasets", [])[:3]) if isinstance(summary, dict) else ""
        rare_metrics = "、".join(summary.get("rare_metrics", [])[:3]) if isinstance(summary, dict) else ""
        gap_overview = str(state.result_schema.get("gap_analysis_overview", "")).strip()
        selected_reason = recommendation.get("selected_reason", selected.recommendation_reason if selected else "")
        lines = [
            "# 创新点分析报告",
            "",
            "## 总体结论",
            f"- 分析模式：{mode}",
            f"- 候选数量：{len(state.innovation_candidates)}",
            f"- 推荐方案：{selected.claim if selected else '尚未生成'}",
            f"- 推荐理由：{selected_reason or '待生成'}",
            f"- 主流方法共性：{common_methods or '待补充'}",
            f"- 主流数据共性：{common_datasets or '待补充'}",
            f"- 主流评测共性：{common_metrics or '待补充'}",
            f"- 常见局限：{common_limitations or '待补充'}",
            f"- 罕见方法线索：{rare_methods or '待补充'}",
            f"- 罕见数据线索：{rare_datasets or '待补充'}",
            f"- 罕见评测线索：{rare_metrics or '待补充'}",
            f"- 最明显 gap：{gap_overview or '待补充'}",
            (
                "- 审核提示：当前结果含 fallback 占位推荐，建议补充文献后再确认最终创新点。"
                if mode == "fallback" or any(item.evidence_mode == "fallback" for item in state.innovation_candidates)
                else "- 审核提示：当前候选主要基于真实结构化文献差异分析生成。"
            ),
            "",
        ]
        if isinstance(summary, dict):
            lines.extend([
                "## 差异分析摘要",
                f"- 方法 gap：{self._format_gap_summary_entries(summary.get('method_gaps', []))}",
                f"- 数据 gap：{self._format_gap_summary_entries(summary.get('data_gaps', []))}",
                f"- 场景 gap：{self._format_gap_summary_entries(summary.get('scenario_gaps', []))}",
                f"- 评价 gap：{self._format_gap_summary_entries(summary.get('evaluation_gaps', []))}",
                f"- 支撑证据映射：{self._format_evidence_map(summary.get('support_evidence_map', {}))}",
                f"- 对照证据映射：{self._format_evidence_map(summary.get('contrast_evidence_map', {}))}",
                "",
            ])
        for idx, item in enumerate(state.innovation_candidates, start=1):
            analysis_basis = "；".join(item.analysis_basis) or "无"
            supporting_evidence = "；".join(item.supporting_evidence) or "无"
            contrast_evidence = "；".join(item.contrast_evidence) or "无"
            lines.extend(
                [
                    f"## 候选 {idx}",
                    f"- 创新点描述：{item.claim}",
                    f"- gap 类型：{item.gap_type}",
                    f"- 支撑文献：{'；'.join(item.supporting_papers) or '无'}",
                    f"- 对照文献：{'；'.join(item.contrast_papers) or '无'}",
                    f"- 分析依据：{analysis_basis}",
                    f"- 支撑证据：{supporting_evidence}",
                    f"- 对照依据：{contrast_evidence}",
                    f"- 创新性说明：{item.novelty_reason}",
                    f"- 少见原因：{item.rare_reason}",
                    f"- 推荐理由：{item.recommendation_reason or '待排序阶段生成'}",
                    f"- 当前名次说明：第 {idx} 名，依据综合评分、证据强度与本科适配度排序。",
                    f"- 证据链解释：优先展示结构化文献中的 limitations / metrics / dataset / problem 摘要，以说明为何判断为该 gap。",
                    (
                        f"- 多维评分：overall={item.overall_score:.2f} / novelty={item.novelty_score:.1f} / "
                        f"feasibility={item.feasibility_score:.1f} / risk={item.risk_score:.1f} / "
                        f"cost={item.experiment_cost:.1f} / undergrad_fit={item.undergrad_fit:.1f} / "
                        f"evidence={item.evidence_strength:.1f}"
                    ),
                    f"- 风险：{item.risk}",
                    f"- 验证计划：{item.verification_plan}",
                    f"- 证据模式：{item.evidence_mode}",
                    (
                        "- 复核建议：该候选为 fallback 占位推荐，进入实验设计前应先补充更多结构化文献证据。"
                        if item.evidence_mode == "fallback"
                        else "- 复核建议：该候选已具备 real 证据，可直接进入实验设计并继续补强实验细节。"
                    ),
                    "",
                ]
            )
        return "\n".join(lines)

    def _format_gap_summary_entries(self, entries: object) -> str:
        if not isinstance(entries, list):
            return "待补充"
        fragments: list[str] = []
        for entry in entries[:2]:
            if not isinstance(entry, dict):
                continue
            focus = str(entry.get("focus", "待补充")).strip() or "待补充"
            description = str(entry.get("description", "")).strip()
            basis = "；".join(str(item) for item in entry.get("basis", [])[:2]) if isinstance(entry.get("basis"), list) else ""
            fragment = f"{focus}：{description}" if description else focus
            if basis:
                fragment = f"{fragment}（依据：{basis}）"
            fragments.append(fragment)
        return "；".join(fragments) or "待补充"

    def _format_evidence_map(self, mapping: object) -> str:
        if not isinstance(mapping, dict):
            return "待补充"
        chunks: list[str] = []
        for label, entries in mapping.items():
            if not isinstance(entries, list) or not entries:
                continue
            first = entries[0]
            if not isinstance(first, dict):
                continue
            phrase = str(first.get("phrase", "待补充")).strip() or "待补充"
            papers = "、".join(first.get("supporting_papers", [])[:2]) if isinstance(first.get("supporting_papers"), list) else ""
            chunks.append(f"{label}={phrase}{f'（{papers}）' if papers else ''}")
        return "；".join(chunks) or "待补充"

    def _experiment_plan_text(self, state: ProjectState) -> str:
        plan = state.experiment_plan
        if not plan:
            return "实验计划尚未生成。"
        return "\n".join(
            [
                "# 实验设计书",
                f"数据集：{', '.join(plan.dataset)}",
                *[f"- 数据说明：{item}" for item in plan.dataset_notes],
                f"基线：{', '.join(plan.baselines)}",
                *[f"- 基线说明：{item}" for item in plan.baseline_notes],
                f"指标：{', '.join(plan.metrics)}",
                *[f"- 指标说明：{item}" for item in plan.metric_notes],
                f"消融：{', '.join(plan.ablations)}",
                f"环境：{', '.join(plan.environment)}",
                "参数：",
                *[f"- {item}" for item in plan.parameters],
                "步骤：",
                *[f"- {step}" for step in plan.steps],
                "运行命令：",
                *[f"- {label}: {command}" for label, command in plan.run_commands.items()],
                "结果文件：",
                *[f"- {item}" for item in plan.result_files],
                "证据来源：",
                *[f"- {item}" for item in plan.evidence_links],
                "预期输出：",
                *[f"- {item}" for item in plan.expected_outputs],
            ]
        )

    def _thesis_text(self, state: ProjectState) -> str:
        if state.paper_document:
            parts = [f"# {state.paper_document.title}", ""]
            for node in state.paper_document.nodes:
                parts.append(f"## {node.title}")
                rendered = self._render_paper_node_markdown(node)
                if rendered:
                    parts.append(rendered)
                if "实验" in node.title:
                    parts.extend(self._format_result_analysis_blocks(state))
                parts.append("")
            return "\n".join(parts)
        manifest = state.template_manifest or TemplateManifest([], {}, [], [], [], "", {}, {})
        template_name = state.template_source.template_name if state.template_source else "未设置"
        parts = [
            "# 毕业论文",
            f"模板来源：{template_name}",
            f"正文结构来源：{' / '.join(state.paper_outline)}",
            "",
        ]
        for section in state.paper_outline:
            parts.append(f"## {section}")
            parts.append(state.paper_sections.get(section, ""))
            if "实验" in section:
                parts.extend(self._format_result_analysis_blocks(state))
            parts.append("")
        return "\n".join(parts)

    def _render_paper_node_markdown(self, node: PaperNode) -> str:
        parts: list[str] = []
        for paragraph in node.paragraphs:
            if paragraph.strip():
                parts.append(paragraph.strip())
        for child in node.children:
            parts.append(f"### {child.title}")
            parts.extend(item.strip() for item in child.paragraphs if item.strip())
        return "\n\n".join(parts).strip()

    def _ppt_text(self, state: ProjectState) -> str:
        lines = ["答辩 PPT 结构", ""]
        draft_mode = getattr(state.request, "delivery_mode", "draft") != "final"
        ppt_summary = str(state.result_schema.get("result_summary_for_ppt", "")).strip()
        mapping = state.result_schema.get("ppt_section_mapping", {})
        result_tables = state.result_schema.get("result_tables", [])
        result_figures = state.result_schema.get("result_figures", [])
        for idx, slide in enumerate(state.ppt_outline, start=1):
            lines.append(f"{idx}. {slide}")
            if slide == "方法设计" and state.selected_innovation:
                lines.append(f"- 方法摘要：{state.selected_innovation.claim}")
            if slide == "实验设置" and state.experiment_plan:
                lines.append(f"- 数据集：{'、'.join(state.experiment_plan.dataset)}")
                lines.append(f"- 指标：{'、'.join(state.experiment_plan.metrics)}")
            if slide == "结果分析" and ppt_summary:
                label = "结果回填说明" if draft_mode else "结果摘要"
                lines.append(f"- {label}：{ppt_summary}")
                if draft_mode:
                    lines.append("- 当前为初稿模式，结果页仅说明待回填的结果表、图表与分析要点。")
                else:
                    lines.extend(self._format_result_slide_details(result_tables, result_figures))
            if isinstance(mapping, dict) and slide in mapping:
                lines.append(f"- 章节映射：{mapping[slide]}")
        return "\n".join(lines)

    def _format_result_analysis_blocks(self, state: ProjectState) -> list[str]:
        blocks: list[str] = []
        draft_mode = getattr(state.request, "delivery_mode", "draft") != "final"
        analysis_text = str(state.result_schema.get("result_analysis_text", "")).strip()
        if analysis_text:
            blocks.append("### 结果分析摘要")
            blocks.append(analysis_text)
        result_tables = state.result_schema.get("result_tables", [])
        if isinstance(result_tables, list) and not draft_mode:
            for table in result_tables[:2]:
                if not isinstance(table, dict):
                    continue
                title = str(table.get("title") or table.get("name") or "结果表").strip()
                summary = str(table.get("summary", "")).strip()
                columns = table.get("columns", [])
                rows = table.get("rows", [])
                blocks.append(f"### {title}")
                if summary:
                    blocks.append(summary)
                if isinstance(columns, list) and columns:
                    blocks.append(f"- 字段：{'、'.join(str(item) for item in columns)}")
                if isinstance(rows, list) and rows:
                    first_row = rows[0]
                    if isinstance(first_row, dict):
                        row_text = "；".join(f"{key}={value}" for key, value in first_row.items())
                        blocks.append(f"- 示例行：{row_text}")
        result_figures = state.result_schema.get("result_figures", [])
        if isinstance(result_figures, list) and not draft_mode:
            for figure in result_figures[:2]:
                if not isinstance(figure, dict):
                    continue
                title = str(figure.get("title") or figure.get("name") or "结果图").strip()
                caption = str(figure.get("caption", "")).strip()
                insight = str(figure.get("insight", "")).strip()
                blocks.append(f"### {title}")
                if caption:
                    blocks.append(f"- 图表说明：{caption}")
                if insight:
                    blocks.append(f"- 分析结论：{insight}")
        if draft_mode and state.experiment_plan:
            blocks.append("### 结果回填说明")
            blocks.append(
                "当前版本为初稿论文，需在完成本地实验后回填主结果对比、消融实验、图表与误差分析。"
            )
            blocks.append(f"- 结果文件：{'、'.join(state.experiment_plan.result_files)}")
        return blocks

    def _format_result_slide_details(self, result_tables: object, result_figures: object) -> list[str]:
        lines: list[str] = []
        if isinstance(result_tables, list):
            for table in result_tables[:2]:
                if not isinstance(table, dict):
                    continue
                title = str(table.get("title") or table.get("name") or "结果表").strip()
                summary = str(table.get("summary", "")).strip()
                if title:
                    lines.append(f"- 表格：{title}")
                if summary:
                    lines.append(f"- 结论：{summary}")
        if isinstance(result_figures, list):
            for figure in result_figures[:2]:
                if not isinstance(figure, dict):
                    continue
                title = str(figure.get("title") or figure.get("name") or "结果图").strip()
                insight = str(figure.get("insight", "")).strip()
                if title:
                    lines.append(f"- 图表：{title}")
                if insight:
                    lines.append(f"- 图示结论：{insight}")
        return lines

    def _write_text(self, path: Path, content: str) -> str:
        path.write_text(content, encoding="utf-8")
        return str(path)

    def _write_docx_like(
        self,
        path: Path,
        content: str,
        *,
        template_path: str | None = None,
        manifest: TemplateManifest | None = None,
        placeholder_values: dict[str, str] | None = None,
    ) -> str:
        try:
            from docx import Document  # type: ignore

            document = self._load_docx_document(Document, template_path)
            if placeholder_values and self._replace_docx_placeholders(document, placeholder_values):
                document.save(path)
                return str(path)
            for line in content.splitlines():
                self._append_docx_line(document, line, manifest)
            document.save(path)
            return str(path)
        except Exception:
            fallback = path.with_suffix(".md")
            fallback.write_text(content, encoding="utf-8")
            return str(fallback)

    def _write_thesis_docx(
        self,
        path: Path,
        state: ProjectState,
        *,
        template_path: str | None = None,
        manifest: TemplateManifest | None = None,
    ) -> str:
        html_path = path.with_name("thesis_source.html")
        if not html_path.exists():
            html_path.write_text(self._render_thesis_html(state), encoding="utf-8")
        html_first_enabled = not template_path and not manifest
        placeholder_values = self._thesis_placeholder_values(state)
        if manifest:
            for section in manifest.section_mapping:
                if section and f"section.{section}" not in placeholder_values:
                    self._resolve_section_text(section, state)
        if not state.paper_document:
            if self._convert_html_to_docx(html_path, path):
                if html_first_enabled:
                    return str(path)
                path.unlink(missing_ok=True)
            result = self._write_docx_like(
                path,
                self._thesis_text(state),
                template_path=template_path,
                manifest=manifest,
                placeholder_values=placeholder_values,
            )
            if Path(result).suffix.lower() == ".docx":
                return result
            warning = "论文 Word 渲染失败，已降级为最小 DOCX 文档。"
            state.warnings.append(warning)
            return self._write_minimal_docx_fallback(path, self._thesis_text(state), warning)
        try:
            if html_first_enabled and self._convert_html_to_docx(html_path, path):
                return str(path)
            from docx import Document  # type: ignore

            document = self._load_docx_document(Document, template_path)
            self._replace_docx_placeholders(document, self._cover_placeholder_values(placeholder_values))
            self._clear_section_placeholders(document)
            self._render_paper_document(document, state, manifest)
            document.save(path)
            return str(path)
        except Exception as exc:
            warning = f"论文 Word 渲染失败，已降级为最小 DOCX 文档：{exc}"
            state.warnings.append(warning)
            return self._write_minimal_docx_fallback(path, self._thesis_text(state), warning)

    def _write_minimal_docx_fallback(self, path: Path, content: str, warning: str) -> str:
        from docx import Document  # type: ignore

        document = Document()
        document.add_heading("论文导出降级结果", level=1)
        document.add_paragraph(warning)
        for block in [segment.strip() for segment in content.split("\n\n") if segment.strip()]:
            document.add_paragraph(block)
        document.save(path)
        return str(path)

    def _load_docx_document(self, document_cls, template_path: str | None):
        if template_path:
            candidate = Path(template_path)
            if candidate.exists():
                try:
                    return document_cls(str(candidate))
                except Exception:
                    pass
        return document_cls()

    def _cover_placeholder_values(self, placeholder_values: dict[str, str]) -> dict[str, str]:
        return {
            key: value
            for key, value in placeholder_values.items()
            if key.startswith("cover.") or key.startswith("thesis.")
        }

    def _thesis_placeholder_values(self, state: ProjectState) -> dict[str, str]:
        manifest = state.template_manifest or TemplateManifest([], {}, [], [], [], "", {}, {})
        template_name = state.template_source.template_name if state.template_source else "未设置"
        values: dict[str, str] = {
            "cover.题目": state.request.topic,
            "thesis.title": state.request.topic,
            "thesis.template_name": template_name,
        }
        for field_name in manifest.cover_fields:
            values.setdefault(f"cover.{field_name}", "")
        section_names = list(dict.fromkeys([*state.paper_outline, *state.paper_sections.keys()]))
        for section in section_names:
            section_text = self._resolve_section_text(section, state)
            values[f"section.{section}"] = section_text
            for alias in self._section_aliases(section):
                values.setdefault(f"section.{alias}", section_text)
        return values

    def _resolve_section_text(self, section: str, state: ProjectState) -> str:
        direct = str(state.paper_sections.get(section, "")).strip()
        if direct:
            return direct
        for alias in self._section_aliases(section):
            alias_text = str(state.paper_sections.get(alias, "")).strip()
            if alias_text:
                if f"论文模板章节“{section}”使用别名“{alias}”填充。" not in state.warnings:
                    state.warnings.append(f"论文模板章节“{section}”使用别名“{alias}”填充。")
                return alias_text
        if f"论文模板章节“{section}”缺少正文内容，已写入待补充提示。" not in state.warnings:
            state.warnings.append(f"论文模板章节“{section}”缺少正文内容，已写入待补充提示。")
        return f"【待补充：{section} 内容尚未生成，请检查章节映射或上游结果摘要。】"

    def _section_aliases(self, section: str) -> list[str]:
        aliases = [section]
        stripped = re.sub(r"^第[0-9一二三四五六七八九十]+章\s*", "", section).strip()
        if stripped and stripped not in aliases:
            aliases.append(stripped)
        for group in SECTION_ALIAS_GROUPS:
            if section in group or stripped in group:
                for candidate in group:
                    if candidate not in aliases:
                        aliases.append(candidate)
        return aliases

    def _replace_docx_placeholders(self, document, placeholder_values: dict[str, str]) -> int:
        replaced = 0
        for paragraph in list(document.paragraphs):
            replaced += self._replace_docx_placeholders_in_paragraph(paragraph, placeholder_values)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in list(cell.paragraphs):
                        replaced += self._replace_docx_placeholders_in_paragraph(paragraph, placeholder_values)
        return replaced

    def _clear_section_placeholders(self, document) -> None:
        for paragraph in list(document.paragraphs):
            original = paragraph.text
            if not original or "{{section." not in original:
                continue
            updated = DOCX_PLACEHOLDER_PATTERN.sub(
                lambda match: "" if match.group(1).startswith("section.") else match.group(0),
                original,
            )
            self._update_paragraph_runs(paragraph, updated)
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in list(cell.paragraphs):
                        original = paragraph.text
                        if not original or "{{section." not in original:
                            continue
                        updated = DOCX_PLACEHOLDER_PATTERN.sub(
                            lambda match: "" if match.group(1).startswith("section.") else match.group(0),
                            original,
                        )
                        self._update_paragraph_runs(paragraph, updated)

    def _render_paper_document(self, document, state: ProjectState, manifest: TemplateManifest | None) -> None:
        paper_document = state.paper_document or self._paper_document_from_sections(state)
        if not paper_document:
            for line in self._thesis_text(state).splitlines():
                self._append_docx_line(document, line, manifest)
            return

        self._append_document_heading(document, paper_document.title, manifest, level=0)
        self._append_body_paragraph(
            document,
            f"模板来源：{state.template_source.template_name if state.template_source else '未设置'}",
            manifest,
        )
        self._append_toc_placeholder(document, paper_document, manifest)
        for node in paper_document.nodes:
            self._render_paper_node(document, node, state, manifest)

    def _paper_document_from_sections(self, state: ProjectState) -> PaperDocument | None:
        if not state.paper_sections:
            return None
        nodes = [
            PaperNode(title=title, level=1, paragraphs=[content] if content else [])
            for title, content in state.paper_sections.items()
        ]
        return PaperDocument(title=state.request.topic or "毕业论文", nodes=nodes)

    def _append_toc_placeholder(self, document, paper_document: PaperDocument, manifest: TemplateManifest | None) -> None:
        self._add_docx_paragraph(document, "目录", self._resolve_docx_style(document, [self._manifest_style(manifest, "chapter"), "Heading 1"]), heading_level=1)
        for index, node in enumerate(paper_document.nodes, start=1):
            self._append_body_paragraph(document, f"{index}. {node.title}", manifest)
            for child_index, child in enumerate(node.children, start=1):
                self._append_body_paragraph(document, f"{index}.{child_index} {child.title}", manifest)

    def _render_paper_node(self, document, node: PaperNode, state: ProjectState, manifest: TemplateManifest | None) -> None:
        if node.title not in {"封面", "目录"}:
            document.add_page_break()
        self._append_document_heading(document, node.title, manifest, level=1)
        for paragraph in node.paragraphs:
            self._append_body_paragraph(document, paragraph, manifest)
        for child in node.children:
            self._append_document_heading(document, child.title, manifest, level=min(child.level, 3))
            for paragraph in child.paragraphs:
                self._append_body_paragraph(document, paragraph, manifest)
            if "主结果" in child.title or "消融" in child.title:
                self._render_result_tables(document, state, manifest)
                self._render_result_figures(document, state, manifest)
        if "参考文献" in node.title:
            self._render_reference_list(document, state, manifest)

    def _append_document_heading(self, document, text: str, manifest: TemplateManifest | None, *, level: int) -> None:
        style_map = {0: "title", 1: "chapter", 2: "section", 3: "subsection"}
        style_name = self._resolve_docx_style(
            document,
            [self._manifest_style(manifest, style_map.get(level, "section")), f"Heading {max(1, min(level, 3))}"] if level else [self._manifest_style(manifest, "title"), "Title"],
        )
        heading_level = None if level == 0 else max(1, min(level, 3))
        self._add_docx_paragraph(document, text, style_name, heading_level=heading_level)

    def _append_body_paragraph(self, document, text: str, manifest: TemplateManifest | None) -> None:
        style_name = self._resolve_docx_style(document, [self._manifest_style(manifest, "body"), "Normal"])
        self._add_docx_paragraph(document, text, style_name)

    def _render_result_tables(self, document, state: ProjectState, manifest: TemplateManifest | None) -> None:
        result_tables = state.result_schema.get("result_tables", [])
        if not isinstance(result_tables, list):
            return
        for table in result_tables[:2]:
            if not isinstance(table, dict):
                continue
            title = str(table.get("title") or table.get("name") or "结果表").strip()
            if title:
                self._append_body_paragraph(document, title, manifest)
            columns = table.get("columns", [])
            rows = table.get("rows", [])
            if not isinstance(columns, list) or not columns:
                continue
            doc_table = document.add_table(rows=1, cols=len(columns))
            doc_table.style = "Table Grid"
            for idx, column in enumerate(columns):
                doc_table.rows[0].cells[idx].text = str(column)
            for row in rows[:3] if isinstance(rows, list) else []:
                if not isinstance(row, dict):
                    continue
                cells = doc_table.add_row().cells
                for idx, column in enumerate(columns):
                    cells[idx].text = str(row.get(column, ""))
            summary = str(table.get("summary", "")).strip()
            if summary:
                self._append_body_paragraph(document, summary, manifest)

    def _render_result_figures(self, document, state: ProjectState, manifest: TemplateManifest | None) -> None:
        result_figures = state.result_schema.get("result_figures", [])
        if not isinstance(result_figures, list):
            return
        for figure in result_figures[:2]:
            if not isinstance(figure, dict):
                continue
            title = str(figure.get("title") or figure.get("name") or "结果图").strip()
            caption = str(figure.get("caption", "")).strip()
            insight = str(figure.get("insight", "")).strip()
            if title:
                self._append_body_paragraph(document, title, manifest)
            if caption:
                self._append_body_paragraph(document, f"图表说明：{caption}", manifest)
            if insight:
                self._append_body_paragraph(document, f"分析结论：{insight}", manifest)

    def _render_reference_list(self, document, state: ProjectState, manifest: TemplateManifest | None) -> None:
        references = [
            f"[{idx + 1}] {record.authors}. {record.title}. {record.year}."
            for idx, record in enumerate(state.literature_records[:8])
            if record.title
        ]
        if not references:
            self._append_body_paragraph(document, "参考文献待根据文献检索结果补充。", manifest)
            return
        for item in references:
            self._append_body_paragraph(document, item, manifest)

    def _replace_docx_placeholders_in_paragraph(self, paragraph, placeholder_values: dict[str, str]) -> int:
        original = paragraph.text
        if not original or "{{" not in original or "}}" not in original:
            return 0

        replaced_count = 0

        def replace_match(match: re.Match[str]) -> str:
            nonlocal replaced_count
            key = match.group(1)
            if key not in placeholder_values:
                replaced_count += 1
                return f"【待补充：{key}】"
            replaced_count += 1
            return placeholder_values[key]

        updated = DOCX_PLACEHOLDER_PATTERN.sub(replace_match, original)
        if replaced_count:
            self._update_paragraph_runs(paragraph, updated)
        return replaced_count

    def _update_paragraph_runs(self, paragraph, updated: str) -> None:
        runs = list(paragraph.runs)
        if not runs:
            paragraph.add_run(updated)
            return
        first_run = runs[0]
        first_run.text = updated
        for run in runs[1:]:
            run.text = ""

    def _append_docx_line(self, document, line: str, manifest: TemplateManifest | None) -> None:
        if line.startswith("# "):
            style_name = self._resolve_docx_style(
                document,
                [
                    self._manifest_style(manifest, "title"),
                    "Title",
                    self._manifest_style(manifest, "chapter"),
                    "Heading 1",
                ],
            )
            self._add_docx_paragraph(document, line[2:], style_name, heading_level=1)
            return

        if line.startswith("## "):
            style_name = self._resolve_docx_style(
                document,
                [
                    self._manifest_style(manifest, "chapter"),
                    "Heading 1",
                    self._manifest_style(manifest, "section"),
                    "Heading 2",
                ],
            )
            self._add_docx_paragraph(document, line[3:], style_name, heading_level=2)
            return

        style_name = self._resolve_docx_style(
            document,
            [self._manifest_style(manifest, "body"), "Normal"],
        )
        self._add_docx_paragraph(document, line, style_name)

    def _manifest_style(self, manifest: TemplateManifest | None, key: str) -> str | None:
        if not manifest:
            return None
        return manifest.style_mapping.get(key)

    def _resolve_docx_style(self, document, candidates: list[str | None]) -> str | None:
        available: dict[str, str] = {}
        for style in document.styles:
            name = getattr(style, "name", None)
            style_id = getattr(style, "style_id", None)
            for value in (name, style_id):
                normalized = self._normalize_docx_style_key(value)
                if normalized and normalized not in available:
                    available[normalized] = name or style_id
        for candidate in candidates:
            normalized = self._normalize_docx_style_key(candidate)
            if normalized and normalized in available:
                return available[normalized]
        return None

    def _normalize_docx_style_key(self, value: str | None) -> str:
        if not value:
            return ""
        collapsed = re.sub(r"[\s_-]+", "", str(value).strip().lower())
        alias_map = {
            "heading1": "heading1",
            "heading2": "heading2",
            "heading3": "heading3",
            "title": "title",
            "normal": "normal",
            "bodytext": "bodytext",
            "bodytextfirstindent": "bodytextfirstindent",
            "caption": "caption",
            "tablegrid": "tablegrid",
        }
        return alias_map.get(collapsed, collapsed)

    def _add_docx_paragraph(
        self,
        document,
        text: str,
        style_name: str | None,
        *,
        heading_level: int | None = None,
    ) -> None:
        if style_name:
            document.add_paragraph(text, style=style_name)
            return
        if heading_level is not None:
            self._add_manual_heading(document, text, heading_level)
            return
        document.add_paragraph(text)

    def _add_manual_heading(self, document, text: str, heading_level: int) -> None:
        paragraph = document.add_paragraph()
        run = paragraph.add_run(text)
        run.bold = True

        if heading_level <= 1:
            run.font.size = Pt(16)
            paragraph.paragraph_format.space_before = Pt(18)
            paragraph.paragraph_format.space_after = Pt(12)
            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER if len(text) <= 24 and not text.startswith("第") else WD_ALIGN_PARAGRAPH.LEFT
        elif heading_level == 2:
            run.font.size = Pt(14)
            paragraph.paragraph_format.space_before = Pt(12)
            paragraph.paragraph_format.space_after = Pt(6)
        else:
            run.font.size = Pt(12)
            paragraph.paragraph_format.space_before = Pt(6)
            paragraph.paragraph_format.space_after = Pt(3)

    def _write_minimal_pdf(self, path: Path, content: str) -> str:
        wrapped = textwrap.wrap(
            content[:1800].replace("(", "[").replace(")", "]"),
            width=80,
        ) or ["thesis assistant"]
        text_commands = " T* ".join([f"({line}) Tj" for line in wrapped[:25]])
        stream = f"BT /F1 12 Tf 50 780 Td 14 TL {text_commands} ET"
        pdf = (
            "%PDF-1.4\n"
            "1 0 obj<< /Type /Catalog /Pages 2 0 R >>endobj\n"
            "2 0 obj<< /Type /Pages /Kids [3 0 R] /Count 1 >>endobj\n"
            "3 0 obj<< /Type /Page /Parent 2 0 R /MediaBox [0 0 595 842] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>endobj\n"
            f"4 0 obj<< /Length {len(stream)} >>stream\n{stream}\nendstream endobj\n"
            "5 0 obj<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>endobj\n"
            "xref\n0 6\n0000000000 65535 f \n"
            "0000000010 00000 n \n0000000063 00000 n \n0000000122 00000 n \n0000000250 00000 n \n0000000000 00000 n \n"
            "trailer<< /Size 6 /Root 1 0 R >>\nstartxref\n0\n%%EOF"
        )
        path.write_bytes(pdf.encode("latin-1", errors="ignore"))
        return str(path)

    def _write_code_zip(self, path: Path, state: ProjectState) -> str:
        with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as archive:
            for relative_path, content in state.generated_code_files.items():
                archive.writestr(relative_path, content)
        return str(path)

    def _write_pptx_like(
        self,
        path: Path,
        content: str,
        *,
        template_path: str | None = None,
        manifest: TemplateManifest | None = None,
    ) -> str:
        try:
            from pptx import Presentation  # type: ignore
            from pptx.util import Inches  # type: ignore

            presentation = self._load_presentation(Presentation, template_path)
            slides = _split_slides(content, manifest.ppt_layouts if manifest else [])
            for slide_text in slides:
                slide_layout = self._choose_slide_layout(presentation, slide_text["layout_hint"])
                slide = presentation.slides.add_slide(slide_layout)
                if slide.shapes.title is not None:
                    slide.shapes.title.text = slide_text["title"]
                body_placeholder = self._find_body_placeholder(slide)
                if body_placeholder is not None:
                    body_placeholder.text = slide_text["body"]
                else:
                    text_box = slide.shapes.add_textbox(
                        Inches(1.0),
                        Inches(1.8),
                        Inches(8.0),
                        Inches(4.5),
                    )
                    text_box.text_frame.text = slide_text["body"]
            presentation.save(path)
            return str(path)
        except Exception:
            fallback = path.with_suffix(".md")
            fallback.write_text(content, encoding="utf-8")
            return str(fallback)

    def _load_presentation(self, presentation_cls, template_path: str | None):
        if template_path:
            candidate = Path(template_path)
            if candidate.exists():
                try:
                    return presentation_cls(str(candidate))
                except Exception:
                    pass
        return presentation_cls()

    def _choose_slide_layout(self, presentation, layout_hint: str):
        normalized_hint = layout_hint.strip().lower()
        if normalized_hint:
            for layout in presentation.slide_layouts:
                layout_name = getattr(layout, "name", "") or ""
                if normalized_hint in layout_name.lower() or layout_name.lower() in normalized_hint:
                    return layout
        if len(presentation.slide_layouts) > 1:
            return presentation.slide_layouts[1]
        return presentation.slide_layouts[0]

    def _find_body_placeholder(self, slide):
        title = slide.shapes.title
        for shape in slide.placeholders:
            if shape is title:
                continue
            if hasattr(shape, "text_frame"):
                return shape
        return None


def _split_slides(content: str, layout_hints: list[str]) -> list[dict[str, str]]:
    slides: list[dict[str, str]] = []
    current_title = "答辩概览"
    current_lines: list[str] = []
    slide_index = 0
    for raw in content.splitlines():
        line = raw.strip()
        if not line:
            continue
        match = re.match(r"^\d+\.\s+(.*)$", line)
        if match:
            if current_lines:
                slides.append(
                    {
                        "title": current_title,
                        "body": "\n".join(current_lines),
                        "layout_hint": _resolve_layout_hint(layout_hints, slide_index),
                    }
                )
                slide_index += 1
            current_title = match.group(1)
            current_lines = [f"{current_title} 的关键信息概览。"]
        else:
            current_lines.append(line)
    if current_lines:
        slides.append(
            {
                "title": current_title,
                "body": "\n".join(current_lines),
                "layout_hint": _resolve_layout_hint(layout_hints, slide_index),
            }
        )
    return slides or [{"title": "答辩概览", "body": content, "layout_hint": ""}]


def _resolve_layout_hint(layout_hints: list[str], index: int) -> str:
    if not layout_hints:
        return ""
    if index < len(layout_hints):
        return layout_hints[index]
    return layout_hints[-1]
