from __future__ import annotations

import shutil
from unittest import mock
import unittest
import uuid
from pathlib import Path

from app.artifact_service import ArtifactService, DEFAULT_THESIS_HTML_STYLE_PROFILE
from app.domain import InnovationCandidate, LiteratureRecord, PaperDocument, PaperNode, ProjectCreate, ProjectState, TemplateManifest, TemplateSource
from app.storage import ProjectStorage
from app.template_service import TemplateService

try:
    from docx import Document  # type: ignore
except Exception:  # pragma: no cover - environment-dependent
    Document = None

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover - environment-dependent
    Presentation = None


class ArtifactServiceTest(unittest.TestCase):
    def setUp(self) -> None:
        self.root = Path("tests_runtime") / f"artifacts-{uuid.uuid4().hex[:8]}"
        self.root.mkdir(parents=True, exist_ok=True)
        self.storage = ProjectStorage(self.root / "projects")
        self.service = ArtifactService(self.storage)
        self.manifest = TemplateManifest(
            section_mapping=["封面", "第1章 绪论"],
            style_mapping={
                "title": "Title",
                "chapter": "Heading 1",
                "section": "Heading 2",
                "body": "Normal",
            },
            cover_fields=["学校"],
            figure_slots=["图1"],
            table_slots=["表1"],
            citation_style="GB/T 7714",
            header_footer_rules={"header": "测试模板", "footer": "页码"},
            toc_rules={"enabled": True, "depth": 3},
            ppt_layouts=["title", "content", "summary"],
        )

    def tearDown(self) -> None:
        shutil.rmtree(self.root, ignore_errors=True)

    def test_default_thesis_html_style_profile_matches_expected_school_rules(self) -> None:
        profile = DEFAULT_THESIS_HTML_STYLE_PROFILE

        self.assertEqual(profile.title.font_family, "宋体")
        self.assertEqual(profile.title.font_size_pt, 16)
        self.assertTrue(profile.title.bold)
        self.assertEqual(profile.title.text_align, "center")

        self.assertEqual(profile.author_meta.font_family, "宋体")
        self.assertEqual(profile.author_meta.font_size_pt, 16)
        self.assertEqual(profile.author_meta.text_align, "center")

        self.assertEqual(profile.cover_label.font_family, "宋体")
        self.assertEqual(profile.cover_label.font_size_pt, 12)
        self.assertEqual(profile.cover_label.text_align, "left")

        self.assertEqual(profile.heading_1.font_family, "宋体")
        self.assertEqual(profile.heading_1.font_size_pt, 15)
        self.assertTrue(profile.heading_1.bold)

        self.assertEqual(profile.heading_2.font_family, "宋体")
        self.assertEqual(profile.heading_2.font_size_pt, 14)
        self.assertTrue(profile.heading_2.bold)

        self.assertEqual(profile.body.font_family, "宋体")
        self.assertEqual(profile.body.font_size_pt, 12)
        self.assertEqual(profile.body.first_line_indent_chars, 2)

    def test_render_thesis_html_outputs_inline_styles_for_title_headings_and_body(self) -> None:
        state = self._build_state("", "")
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[
                PaperNode(
                    title="第1章 绪论",
                    level=1,
                    children=[
                        PaperNode(title="研究背景", level=2, paragraphs=["这是第一段正文。\n\n这是第二段正文。"]),
                    ],
                )
            ],
        )

        html = self.service._render_thesis_html(state)

        self.assertIn("font-family:宋体", html)
        self.assertIn("font-size:16pt", html)
        self.assertIn("font-size:15pt", html)
        self.assertIn("font-size:14pt", html)
        self.assertIn("font-size:12pt", html)
        self.assertIn("text-align:center", html)
        self.assertIn("text-indent:2em", html)
        self.assertIn("中文文本分类算法", html)
        self.assertIn("第1章 绪论", html)
        self.assertIn("研究背景", html)
        self.assertIn("这是第一段正文。", html)

    def test_render_thesis_html_outputs_template_cover_fields_with_placeholders(self) -> None:
        state = self._build_state("", "")
        state.template_manifest = TemplateManifest(
            section_mapping=["第1章 绪论"],
            style_mapping={},
            cover_fields=["学校", "学院", "专业", "学号", "指导教师"],
            figure_slots=[],
            table_slots=[],
            citation_style="GB/T 7714",
            header_footer_rules={},
            toc_rules={},
        )
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[PaperNode(title="第1章 绪论", level=1, paragraphs=["正文段落。"])],
        )

        html = self.service._render_thesis_html(state)

        self.assertIn("课题名称：中文文本分类算法", html)
        self.assertIn("论文模板：工科毕业论文", html)
        self.assertIn("学校：待填写", html)
        self.assertIn("学院：待填写", html)
        self.assertIn("指导教师：待填写", html)
        self.assertIn("width:720px", html)

    def test_render_thesis_html_includes_result_analysis_blocks_for_experiment_sections(self) -> None:
        state = self._build_state("", "")
        state.request.delivery_mode = "draft"
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[PaperNode(title="实验", level=1, paragraphs=["实验部分概述"])],
        )
        state.result_schema["result_analysis_text"] = "当前处于初稿模式，真实结果需由用户完成实验后回填。"

        html = self.service._render_thesis_html(state)

        self.assertIn("结果分析摘要", html)
        self.assertIn("用户完成实验后回填", html)

    def test_render_all_writes_thesis_html_debug_artifact(self) -> None:
        state = self._build_state("", "")
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[
                PaperNode(
                    title="第1章 绪论",
                    level=1,
                    children=[PaperNode(title="研究背景", level=2, paragraphs=["正文段落。"])],
                )
            ],
        )

        result = self.service.render_all(state)

        html_path = Path(result.artifacts.thesis_html or "")
        self.assertTrue(html_path.exists())
        self.assertEqual(html_path.suffix, ".html")
        html = html_path.read_text(encoding="utf-8")
        self.assertIn("中文文本分类算法", html)
        self.assertIn("font-family:宋体", html)
        self.assertIn("text-indent:2em", html)

    def test_convert_html_to_docx_returns_false_when_pandoc_is_missing(self) -> None:
        html_path = self.root / "sample.html"
        docx_path = self.root / "sample.docx"
        html_path.write_text("<html><body><p>demo</p></body></html>", encoding="utf-8")

        with mock.patch("app.artifact_service.shutil.which", return_value=None):
            converted = self.service._convert_html_to_docx(html_path, docx_path)

        self.assertFalse(converted)
        self.assertFalse(docx_path.exists())

    def test_convert_html_to_docx_invokes_pandoc_when_available(self) -> None:
        html_path = self.root / "sample.html"
        docx_path = self.root / "sample.docx"
        html_path.write_text("<html><body><p>demo</p></body></html>", encoding="utf-8")

        def fake_run(command, capture_output, text, check):
            docx_path.write_bytes(b"docx")

            class Result:
                returncode = 0

            return Result()

        with mock.patch("app.artifact_service.shutil.which", return_value="pandoc"), mock.patch(
            "app.artifact_service.subprocess.run",
            side_effect=fake_run,
        ) as run_mock:
            converted = self.service._convert_html_to_docx(html_path, docx_path)

        self.assertTrue(converted)
        self.assertTrue(docx_path.exists())
        run_mock.assert_called_once()

    def test_write_thesis_docx_prefers_html_conversion_when_available(self) -> None:
        state = self._build_state("", "")
        state.paper_document = None
        target = self.root / "thesis.docx"

        def fake_convert(html_path: Path, docx_path: Path) -> bool:
            self.assertTrue(html_path.exists())
            docx_path.write_bytes(b"html-docx")
            return True

        with mock.patch.object(self.service, "_convert_html_to_docx", side_effect=fake_convert) as convert_mock, mock.patch.object(
            self.service,
            "_write_docx_like",
            side_effect=AssertionError("legacy docx writer should not be used when html conversion succeeds"),
        ):
            result = self.service._write_thesis_docx(target, state)

        self.assertEqual(result, str(target))
        self.assertTrue(target.exists())
        convert_mock.assert_called_once()

    def test_write_thesis_docx_falls_back_to_legacy_path_when_html_conversion_fails(self) -> None:
        state = self._build_state("", "")
        state.paper_document = None
        target = self.root / "fallback-thesis.docx"

        with mock.patch.object(self.service, "_convert_html_to_docx", return_value=False), mock.patch.object(
            self.service,
            "_write_docx_like",
            return_value=str(target),
        ) as legacy_mock:
            result = self.service._write_thesis_docx(target, state)

        self.assertEqual(result, str(target))
        legacy_mock.assert_called_once()

    @unittest.skipIf(
        Document is None or Presentation is None,
        "python-docx or python-pptx is not installed in the current environment",
    )
    def test_render_all_uses_valid_word_and_ppt_templates(self) -> None:
        word_template = self.root / "template.docx"
        ppt_template = self.root / "template.pptx"
        template = Document()
        template.add_paragraph("{{cover.题目}}", style="Title")
        template.add_page_break()
        template.add_heading("摘要", level=1)
        template.add_paragraph("{{section.摘要}}")
        template.add_heading("第1章 绪论", level=1)
        template.add_paragraph("{{section.第1章 绪论}}")
        template.save(word_template)
        Presentation().save(ppt_template)

        state = self._build_state(str(word_template), str(ppt_template))

        result = self.service.render_all(state)

        self.assertTrue(Path(result.artifacts.thesis_docx or "").exists())
        self.assertTrue(Path(result.artifacts.defense_pptx or "").exists())
        self.assertEqual(Path(result.artifacts.thesis_docx or "").suffix, ".docx")
        self.assertEqual(Path(result.artifacts.defense_pptx or "").suffix, ".pptx")

        thesis = Document(result.artifacts.thesis_docx)
        texts = [paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip()]
        self.assertIn("中文文本分类算法", texts)
        self.assertIn("第1章 绪论", texts)

    @unittest.skipIf(
        Document is None or Presentation is None,
        "python-docx or python-pptx is not installed in the current environment",
    )
    def test_render_all_resolves_section_aliases_and_removes_placeholders(self) -> None:
        word_template = self.root / "alias-template.docx"
        ppt_template = self.root / "alias-template.pptx"
        template = Document()
        template.add_paragraph("{{cover.题目}}", style="Title")
        heading = template.add_heading("第1章 绪论", level=1)
        heading.runs[0].bold = True
        body = template.add_paragraph()
        body_run = body.add_run("{{section.绪论}}")
        body_run.bold = True
        template.save(word_template)
        Presentation().save(ppt_template)

        state = self._build_state(str(word_template), str(ppt_template))
        state.paper_outline = ["引言"]
        state.paper_sections = {"引言": "这里是通过章节别名回填的正文内容。"}

        result = self.service.render_all(state)

        thesis = Document(result.artifacts.thesis_docx)
        texts = [paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip()]
        self.assertIn("这里是通过章节别名回填的正文内容。", texts)
        self.assertTrue(all("{{section." not in text for text in texts))
        body_paragraph = next(paragraph for paragraph in thesis.paragraphs if "这里是通过章节别名回填的正文内容。" in paragraph.text)
        self.assertTrue(body_paragraph.runs)
        self.assertTrue(body_paragraph.runs[0].bold)

    @unittest.skipIf(
        Document is None or Presentation is None,
        "python-docx or python-pptx is not installed in the current environment",
    )
    def test_render_all_writes_fallback_notice_for_missing_sections(self) -> None:
        word_template = self.root / "missing-template.docx"
        ppt_template = self.root / "missing-template.pptx"
        template = Document()
        template.add_heading("第4章 实验结果与分析", level=1)
        template.add_paragraph("{{section.第4章 实验结果与分析}}")
        template.save(word_template)
        Presentation().save(ppt_template)

        state = self._build_state(str(word_template), str(ppt_template))
        state.paper_outline = []
        state.paper_sections = {}

        result = self.service.render_all(state)

        thesis = Document(result.artifacts.thesis_docx)
        thesis_text = "\n".join(paragraph.text for paragraph in thesis.paragraphs)
        self.assertIn("待补充", thesis_text)
        self.assertNotIn("{{section.", thesis_text)
        self.assertTrue(any("缺少正文内容" in warning for warning in state.warnings))

    @unittest.skipIf(
        Document is None or Presentation is None,
        "python-docx or python-pptx is not installed in the current environment",
    )
    def test_render_all_outputs_structured_subsections_for_thesis_docx(self) -> None:
        word_template = self.root / "structured-template.docx"
        ppt_template = self.root / "structured-template.pptx"
        template = Document()
        template.add_paragraph("{{cover.题目}}", style="Title")
        template.save(word_template)
        Presentation().save(ppt_template)

        state = self._build_state(str(word_template), str(ppt_template))
        state.paper_outline = ["第1章 绪论", "第4章 实验结果与分析"]
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[
                PaperNode(
                    title="第1章 绪论",
                    level=1,
                    children=[
                        PaperNode(title="研究背景", level=2, paragraphs=["背景段落一。", "背景段落二。"]),
                        PaperNode(title="问题定义", level=2, paragraphs=["问题定义段落一。", "问题定义段落二。"]),
                    ],
                ),
                PaperNode(
                    title="第4章 实验结果与分析",
                    level=1,
                    children=[
                        PaperNode(title="实验环境", level=2, paragraphs=["环境段落一。", "环境段落二。"]),
                        PaperNode(title="主结果分析", level=2, paragraphs=["结果段落一。", "结果段落二。"]),
                    ],
                ),
            ],
        )
        state.paper_sections = {
            "第1章 绪论": "### 研究背景\n\n背景段落一。\n\n### 问题定义\n\n问题定义段落一。",
            "第4章 实验结果与分析": "### 实验环境\n\n环境段落一。\n\n### 主结果分析\n\n结果段落一。",
        }

        result = self.service.render_all(state)

        thesis = Document(result.artifacts.thesis_docx)
        texts = [paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip()]
        self.assertIn("第1章 绪论", texts)
        self.assertIn("研究背景", texts)
        self.assertIn("主结果分析", texts)

    @unittest.skipIf(
        Document is None or Presentation is None,
        "python-docx or python-pptx is not installed in the current environment",
    )
    def test_render_all_preserves_docx_caption_guidance_in_final_thesis_docx(self) -> None:
        word_template = self.root / "caption-template.docx"
        ppt_template = self.root / "caption-template.pptx"
        template = Document()
        template.add_paragraph("{{cover.题目}}", style="Title")
        template.save(word_template)
        Presentation().save(ppt_template)

        state = self._build_state(str(word_template), str(ppt_template))
        state.paper_outline = ["第4章 实验结果与分析", "参考文献"]
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[
                PaperNode(
                    title="第4章 实验结果与分析",
                    level=1,
                    children=[
                        PaperNode(
                            title="主结果分析",
                            level=2,
                            paragraphs=[
                                "本节围绕主结果对比、指标变化和关键现象进行定量分析。",
                                "表 4-1 主结果对比表。建议将主指标、对比基线和提升幅度整理为独立表题，并在表后追加 1 段结果分析。",
                                "图 4-1 训练曲线。建议在图后单独说明收敛趋势、稳定性变化与代表性现象。",
                            ],
                        )
                    ],
                ),
                PaperNode(
                    title="参考文献",
                    level=1,
                    paragraphs=["参考文献导出时按独立段落组织，不使用项目符号列表。"],
                ),
            ],
        )
        state.paper_sections = {
            "第4章 实验结果与分析": "### 主结果分析\n\n表 4-1 主结果对比表。\n\n图 4-1 训练曲线。",
            "参考文献": "参考文献导出时按独立段落组织，不使用项目符号列表。",
        }

        result = self.service.render_all(state)

        thesis = Document(result.artifacts.thesis_docx)
        texts = [paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip()]
        self.assertIn("表 4-1 主结果对比表。建议将主指标、对比基线和提升幅度整理为独立表题，并在表后追加 1 段结果分析。", texts)
        self.assertIn("图 4-1 训练曲线。建议在图后单独说明收敛趋势、稳定性变化与代表性现象。", texts)
        self.assertIn("参考文献导出时按独立段落组织，不使用项目符号列表。", texts)

    def test_render_all_falls_back_when_template_files_are_invalid(self) -> None:
        word_template = self.root / "invalid-template.docx"
        ppt_template = self.root / "invalid-template.pptx"
        word_template.write_text("invalid docx placeholder", encoding="utf-8")
        ppt_template.write_text("invalid pptx placeholder", encoding="utf-8")

        state = self._build_state(str(word_template), str(ppt_template))

        result = self.service.render_all(state)

        self.assertTrue(Path(result.artifacts.thesis_docx or "").exists())
        self.assertTrue(Path(result.artifacts.defense_pptx or "").exists())

    @unittest.skipIf(Document is None, "python-docx is not installed in the current environment")
    def test_thesis_docx_fallback_still_returns_docx_when_rendering_fails(self) -> None:
        state = self._build_state("", "")
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[PaperNode(title="第1章 绪论", level=1, paragraphs=["绪论段落。"])],
        )

        original_render = self.service._render_paper_document

        def raising_render(*args, **kwargs):
            raise RuntimeError("render failure")

        self.service._render_paper_document = raising_render  # type: ignore[assignment]
        try:
            result = self.service.render_all(state)
        finally:
            self.service._render_paper_document = original_render  # type: ignore[assignment]

        thesis_path = Path(result.artifacts.thesis_docx or "")
        self.assertTrue(thesis_path.exists())
        self.assertEqual(thesis_path.suffix, ".docx")
        thesis = Document(str(thesis_path))
        text = "\n".join(paragraph.text for paragraph in thesis.paragraphs)
        self.assertIn("论文导出降级结果", text)
        self.assertTrue(any("论文 Word 渲染失败" in warning for warning in state.warnings))

    @unittest.skipIf(Document is None, "python-docx is not installed in the current environment")
    def test_render_all_avoids_fallback_when_template_has_no_heading_styles(self) -> None:
        word_template = self.root / "styleless-template.docx"
        template = Document()
        title = template.add_paragraph()
        title.add_run("{{cover.题目}}")
        template.add_paragraph("摘要")
        template.add_paragraph("{{section.摘要}}")
        template.save(word_template)

        state = self._build_state(str(word_template), "")
        state.paper_document = PaperDocument(
            title="中文文本分类算法",
            nodes=[PaperNode(title="第1章 绪论", level=1, paragraphs=["绪论段落。"])]
        )

        result = self.service.render_all(state)

        thesis_path = Path(result.artifacts.thesis_docx or "")
        self.assertTrue(thesis_path.exists())
        thesis = Document(str(thesis_path))
        text = "\n".join(paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip())
        self.assertIn("中文文本分类算法", text)
        self.assertIn("第1章 绪论", text)
        self.assertNotIn("论文导出降级结果", text)
        self.assertFalse(any("已降级为最小 DOCX 文档" in warning for warning in state.warnings))

    @unittest.skipIf(Document is None, "python-docx is not installed in the current environment")
    def test_general_undergraduate_library_template_is_renderable(self) -> None:
        service = TemplateService()
        source, manifest = service.choose_default_template(ProjectCreate(topic="中文文本分类算法", paper_type="algorithm"))

        self.assertTrue(Path(source.template_path or "").exists())
        state = self._build_state(source.template_path or "", "")
        state.template_source = source
        state.template_manifest = manifest
        state.paper_outline = ["摘要", "第1章 绪论", "参考文献"]
        state.paper_sections = {
            "摘要": "这是中文摘要。",
            "第1章 绪论": "这是绪论正文。",
            "参考文献": "[1] Tester. Demo. 2026.",
        }

        result = self.service.render_all(state)

        thesis = Document(result.artifacts.thesis_docx)
        texts = [paragraph.text for paragraph in thesis.paragraphs if paragraph.text.strip()]
        self.assertIn("中文文本分类算法", texts)
        self.assertIn("这是中文摘要。", texts)
        self.assertIn("这是绪论正文。", texts)

    def test_literature_review_contains_structured_evidence_fields(self) -> None:
        state = self._build_state("", "")
        state.literature_records = [
            LiteratureRecord(
                source="semantic_scholar",
                title="A Structured Paper",
                authors="Tester",
                year=2025,
                abstract="This paper proposes a structured method.",
                doi_or_url="https://example.com/paper",
                problem="文本分类问题",
                method="结构化编码方法",
                dataset="THUCNews",
                metrics="Accuracy, F1",
                conclusion="方法有效",
                limitations="样本规模有限",
                evidence_source="abstract",
                confidence_score=0.78,
                citation_count=12,
            )
        ]
        state.survey_table = [
            {
                "title": "A Structured Paper",
                "problem": "文本分类问题",
                "method": "结构化编码方法",
                "dataset": "THUCNews",
                "metrics": "Accuracy, F1",
                "conclusion": "方法有效",
                "limitations": "样本规模有限",
                "source": "semantic_scholar",
                "doi_or_url": "https://example.com/paper",
                "evidence_source": "abstract",
                "confidence": "0.78",
                "evidence_quote": "This paper proposes a structured method.",
                "pdf_path": "",
                "pdf_parse_status": "not_applicable",
                "pdf_parse_message": "",
                "citation_count": "12",
                "is_fallback": "no",
                "needs_review": "no",
                "review_note": "",
            }
        ]

        result = self.service.render_all(state)

        exported = Path(result.artifacts.literature_review or "")
        self.assertTrue(exported.exists())
        content = exported.read_bytes()
        if exported.suffix == ".csv":
            text = content.decode("utf-8-sig")
            self.assertIn("evidence_source", text)
            self.assertIn("confidence", text)
            self.assertIn("is_fallback", text)
            self.assertIn("evidence_quote", text)
            self.assertIn("needs_review", text)

    def test_innovation_report_contains_gap_evidence_and_recommendation_fields(self) -> None:
        state = self._build_state("", "")
        state.result_schema["gap_analysis"] = {
            "mode": "real",
            "common_methods": ["Transformer 编码器", "轻量分类头"],
            "common_datasets": ["THUCNews"],
            "common_metrics": ["Accuracy", "F1"],
            "common_limitations": ["鲁棒性评测不足"],
            "rare_methods": ["对比学习模块"],
            "rare_datasets": ["Fudan"],
            "rare_metrics": ["Robustness"],
            "method_gaps": [
                {
                    "focus": "Transformer 编码器与轻量模块组合改造",
                    "description": "主流方法集中，但轻量协同改造较少。",
                    "basis": ["主流方法集中在 Transformer 编码器。"],
                }
            ],
            "data_gaps": [
                {
                    "focus": "THUCNews 外数据补强",
                    "description": "数据覆盖较窄。",
                    "basis": ["大多数论文集中在 THUCNews。"],
                }
            ],
            "scenario_gaps": [
                {
                    "focus": "低资源与跨域验证",
                    "description": "真实应用场景覆盖不足。",
                    "basis": ["多篇论文提到低资源验证不足。"],
                }
            ],
            "evaluation_gaps": [
                {
                    "focus": "鲁棒性与误差分析",
                    "description": "评价维度仍不完整。",
                    "basis": ["4 篇论文只报告 Accuracy/F1。"],
                }
            ],
            "support_evidence_map": {
                "methods": [{"phrase": "Transformer 编码器", "supporting_papers": ["Paper 1", "Paper 2"]}],
            },
            "contrast_evidence_map": {
                "metrics": [{"phrase": "Robustness", "supporting_papers": ["Paper 3"]}],
            },
        }
        state.result_schema["gap_analysis_overview"] = "最明显的 gap 类型是评价空白，因为主流论文主要集中在常规指标。"
        state.result_schema["innovation_recommendation"] = {
            "selected_reason": "评价补强方案成本低、证据强，适合直接进入实验设计。"
        }
        state.innovation_candidates = [
            InnovationCandidate(
                claim="面向中文文本分类的鲁棒性与可解释评价补强方案",
                supporting_papers=["Paper 1", "Paper 2"],
                contrast_papers=["Paper 3"],
                analysis_basis=["4 篇论文主要只报告 Accuracy/F1，缺少鲁棒性评测。"],
                supporting_evidence=["Paper 1：limitations 指向“缺少鲁棒性评测”"],
                contrast_evidence=["Paper 3：metrics 指向“Accuracy, F1”"],
                novelty_reason="当前工作多关注常规精度指标。",
                feasibility_score=8.2,
                risk="指标设计不清晰会削弱说服力。",
                verification_plan="补充鲁棒性测试与错误案例分析。",
                gap_type="evaluation_gap",
                rare_reason="评价维度扩展往往未被当作主要贡献展开。",
                recommendation_reason="成本低、适合本科毕业论文且证据较强。",
                novelty_score=7.5,
                risk_score=4.5,
                experiment_cost=3.5,
                undergrad_fit=8.8,
                evidence_strength=8.0,
                evidence_mode="real",
                overall_score=8.01,
            )
        ]
        state.selected_innovation = state.innovation_candidates[0]

        result = self.service.render_all(state)

        report_path = Path(result.artifacts.innovation_report or "")
        self.assertTrue(report_path.exists())
        report = report_path.read_text(encoding="utf-8")
        self.assertIn("创新点分析报告", report)
        self.assertIn("gap 类型", report)
        self.assertIn("支撑文献", report)
        self.assertIn("少见原因", report)
        self.assertIn("推荐理由", report)
        self.assertIn("证据模式", report)
        self.assertIn("最明显 gap", report)
        self.assertIn("差异分析摘要", report)
        self.assertIn("支撑证据映射", report)
        self.assertIn("对照证据映射", report)
        self.assertIn("证据链解释", report)
        self.assertIn("复核建议", report)
        self.assertIn("分析依据", report)
        self.assertIn("支撑证据", report)
        self.assertIn("对照依据", report)
        self.assertIn("当前名次说明", report)

    def test_thesis_and_ppt_consume_structured_result_analysis(self) -> None:
        state = self._build_state("", "")
        state.request.delivery_mode = "final"
        state.paper_outline = ["摘要", "实验", "结论"]
        state.paper_sections = {
            "摘要": "摘要内容",
            "实验": "实验部分概述",
            "结论": "结论内容",
        }
        state.ppt_outline = ["方法设计", "实验设置", "结果分析", "结论与展望"]
        state.result_schema.update(
            {
                "result_analysis_text": "主结果优于基线，消融实验验证关键模块有效。",
                "result_summary_for_paper": "论文实验章节复用了主结果和消融结论。",
                "result_summary_for_ppt": "PPT 展示主结果对比、消融和训练曲线。",
                "result_tables": [
                    {
                        "title": "主结果对比表",
                        "columns": ["方法", "Accuracy", "F1"],
                        "rows": [{"方法": "本文方案", "Accuracy": 0.91, "F1": 0.89}],
                        "summary": "本文方案在 Accuracy 和 F1 上均优于基线。",
                    }
                ],
                "result_figures": [
                    {
                        "title": "训练曲线",
                        "caption": "训练后期收敛稳定。",
                        "insight": "模型在有限轮次内完成收敛。",
                    }
                ],
                "ppt_section_mapping": {
                    "结果分析": "实验章节/结果分析段",
                    "结论与展望": "结论章节",
                },
            }
        )

        result = self.service.render_all(state)

        thesis_path = Path(result.artifacts.thesis_docx or "")
        ppt_path = Path(result.artifacts.defense_pptx or "")
        self.assertTrue(thesis_path.exists())
        self.assertTrue(ppt_path.exists())
        if thesis_path.suffix == ".md":
            thesis_text = thesis_path.read_text(encoding="utf-8")
        elif Document is not None:
            thesis_doc = Document(thesis_path)
            thesis_text = "\n".join(paragraph.text for paragraph in thesis_doc.paragraphs)
        else:
            thesis_text = ""
        if ppt_path.suffix == ".txt":
            ppt_text = ppt_path.read_text(encoding="utf-8")
        elif Presentation is not None:
            presentation = Presentation(ppt_path)
            slide_texts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slide_texts.append(text)
            ppt_text = "\n".join(slide_texts)
        else:
            ppt_text = ""
        self.assertIn("结果分析摘要", thesis_text)
        self.assertIn("主结果对比表", thesis_text)
        self.assertIn("训练曲线", thesis_text)
        self.assertIn("表格：主结果对比表", ppt_text)
        self.assertIn("图表：训练曲线", ppt_text)
        self.assertIn("章节映射：实验章节/结果分析段", ppt_text)

    def test_draft_thesis_and_ppt_do_not_emit_fake_result_claims(self) -> None:
        state = self._build_state("", "")
        state.request.delivery_mode = "draft"
        state.paper_outline = ["摘要", "实验", "结论"]
        state.paper_sections = {
            "摘要": "摘要内容",
            "实验": "实验部分概述",
            "结论": "结论内容",
        }
        state.ppt_outline = ["方法设计", "实验设置", "结果分析", "结论与展望"]
        state.result_schema.update(
            {
                "result_analysis_text": "当前处于初稿模式，论文仅保留实验设计、实施流程与结果记录模板。真实结果、图表与定量结论需由用户完成实验后回填。",
                "result_summary_for_paper": "实验结果部分暂以结果记录模板占位，待完成主结果对比、消融实验和误差分析后再补入正式论文。",
                "result_summary_for_ppt": "答辩材料当前只展示实验设置与结果回填说明，真实结果页需在实验完成后更新。",
                "result_tables": [
                    {
                        "title": "实验结果记录模板",
                        "columns": ["实验项", "待填写结果", "说明"],
                        "rows": [{"实验项": "主结果对比", "待填写结果": "待实验后填写", "说明": "填写各方案指标。"}],
                        "summary": "当前为初稿模式，结果表由用户完成实验后回填。",
                    }
                ],
                "result_figures": [
                    {
                        "title": "实验图表规划",
                        "caption": "建议补充训练曲线与主结果对比图。",
                        "insight": "当前仅输出图表规划，不生成虚拟结论。",
                    }
                ],
                "ppt_section_mapping": {
                    "结果分析": "实验章节/结果回填段",
                    "结论与展望": "结论章节",
                },
            }
        )

        result = self.service.render_all(state)

        thesis_path = Path(result.artifacts.thesis_docx or "")
        ppt_path = Path(result.artifacts.defense_pptx or "")
        if thesis_path.suffix == ".md":
            thesis_text = thesis_path.read_text(encoding="utf-8")
        elif Document is not None:
            thesis_text = "\n".join(paragraph.text for paragraph in Document(thesis_path).paragraphs)
        else:
            thesis_text = ""
        if ppt_path.suffix == ".txt":
            ppt_text = ppt_path.read_text(encoding="utf-8")
        elif Presentation is not None:
            presentation = Presentation(ppt_path)
            slide_texts: list[str] = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        slide_texts.append(text)
            ppt_text = "\n".join(slide_texts)
        else:
            ppt_text = ""

        self.assertIn("用户完成实验后回填", thesis_text)
        self.assertIn("结果回填说明", ppt_text)
        self.assertNotIn("表格：主结果对比表", ppt_text)
        self.assertNotIn("训练曲线", thesis_text)

    def _build_state(self, word_template_path: str, ppt_template_path: str) -> ProjectState:
        state = ProjectState(
            project_id="project-001",
            request=ProjectCreate(topic="中文文本分类算法"),
        )
        state.template_source = TemplateSource(
            source_type="library_default",
            template_id="engineering_thesis",
            template_name="工科毕业论文",
            template_path=word_template_path,
            ppt_template_path=ppt_template_path,
        )
        state.template_manifest = self.manifest
        state.paper_outline = ["第1章 绪论"]
        state.paper_sections = {"第1章 绪论": "这是论文正文。"}
        state.ppt_outline = ["研究背景", "方法设计", "实验结果"]
        state.result_schema["result_summary_for_ppt"] = "主结果优于基线，消融验证关键模块有效。"
        state.result_schema["ppt_section_mapping"] = {"方法设计": "方法章节", "实验结果": "实验章节"}
        return state

    def test_experiment_plan_export_contains_milestone_three_fields(self) -> None:
        state = self._build_state("", "")
        state.experiment_plan = __import__("app.domain", fromlist=["ExperimentPlan"]).ExperimentPlan(
            dataset=["THUCNews"],
            baselines=["Transformer"],
            metrics=["Accuracy", "F1"],
            ablations=["去除增强模块"],
            environment=["Python 3.11"],
            parameters=["seed=42"],
            dataset_notes=["使用官方训练/验证划分"],
            baseline_notes=["基线来自主流文献"],
            metric_notes=["Accuracy 与 F1 为核心指标"],
            steps=["训练基线", "运行完整模型"],
            expected_outputs=["主结果表"],
            run_commands={"train": "python train.py --config configs/default.yaml"},
            result_files=["results/eval_metrics.json"],
            evidence_links=["Paper A"],
        )

        result = self.service.render_all(state)

        exported = Path(result.artifacts.experiment_plan or "")
        self.assertTrue(exported.exists())
        if exported.suffix == ".docx" and Document is not None:
            text = "\n".join(paragraph.text for paragraph in Document(exported).paragraphs)
        else:
            text = exported.read_text(encoding="utf-8")
        self.assertIn("运行命令", text)
        self.assertIn("结果文件", text)
        self.assertIn("证据来源", text)

    def test_thesis_placeholder_values_ignore_template_only_sections(self) -> None:
        state = self._build_state("", "")
        state.template_manifest = TemplateManifest(
            section_mapping=["封面", "模板致谢", "模板附录"],
            style_mapping={
                "title": "Title",
                "chapter": "Heading 1",
                "section": "Heading 2",
                "body": "Normal",
            },
            cover_fields=["学校"],
            figure_slots=[],
            table_slots=[],
            citation_style="GB/T 7714",
            header_footer_rules={},
            toc_rules={"enabled": True, "depth": 3},
            ppt_layouts=[],
        )
        state.paper_outline = ["摘要", "实验"]
        state.paper_sections = {
            "摘要": "这是摘要。",
            "实验": "这是实验章节。",
        }

        values = self.service._thesis_placeholder_values(state)

        self.assertIn("section.摘要", values)
        self.assertIn("section.实验", values)
        self.assertNotIn("section.模板致谢", values)
        self.assertNotIn("section.模板附录", values)


if __name__ == "__main__":
    unittest.main()
